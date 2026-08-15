from __future__ import annotations

import base64
import csv
import io
import logging
import mimetypes
import tempfile
import uuid
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

from .settings import settings


logger = logging.getLogger(__name__)

TEXT_SUFFIXES = {
    ".txt", ".md", ".csv", ".tsv", ".json", ".py", ".js", ".ts", ".tsx",
    ".jsx", ".html", ".css", ".xml", ".yaml", ".yml", ".sql", ".r",
}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".gif"}
SUPPORTED_SUFFIXES = TEXT_SUFFIXES | IMAGE_SUFFIXES | {".pdf", ".docx", ".pptx", ".xlsx"}

# Soft targets for student and course uploads. Compression never expands a file
# and never replaces a PDF when text extraction would regress.
_PDF_COMPRESS_MIN_BYTES = 400_000
_PDF_IMAGE_MAX_EDGE = 1600
_PDF_JPEG_QUALITY = 72
_IMAGE_MAX_EDGE = 2048
_IMAGE_JPEG_QUALITY = 82
_MAX_PDF_PAGES = 200
_MAX_OFFICE_ARCHIVE_ENTRIES = 5_000
_MAX_OFFICE_UNCOMPRESSED_BYTES = 100 * 1024 * 1024
_MAX_OFFICE_COMPRESSION_RATIO = 100
_MAX_DOCUMENT_PARAGRAPHS = 50_000
_MAX_PRESENTATION_SLIDES = 300
_MAX_SPREADSHEET_CELLS = 200_000


@dataclass(frozen=True)
class StoredUpload:
    name: str
    path: Path
    mime: str
    size: int
    supported: bool
    extracted_text: str = ""
    storage_key: str | None = None
    storage_provider: str = "local"
    source_id: str | None = None

    @property
    def is_image(self) -> bool:
        suffix = (
            Path(self.name).suffix.lower()
            if self.storage_key
            else self.path.suffix.lower()
        )
        return suffix in IMAGE_SUFFIXES

    def read_bytes(self) -> bytes:
        """Return upload bytes from local path or configured object storage."""
        if self.storage_key and self.storage_provider != "local":
            from backend.persistence.factory import get_file_storage

            return get_file_storage().get_bytes(self.storage_key)
        return self.path.read_bytes()


def _safe_name(name: str) -> str:
    safe = Path(name).name.replace("\x00", "").strip()
    return safe[:180] or "upload"


def _upload_root(thread_id: str) -> Path:
    safe_id = "".join(
        character for character in thread_id if character.isalnum() or character in "-_"
    )
    root = (settings.files_dir / "threads" / safe_id / "uploads").resolve()
    if settings.files_dir not in root.parents:
        raise ValueError("Unsafe chat identifier")
    root.mkdir(parents=True, exist_ok=True)
    return root


def _validate_office_archive(path: Path) -> None:
    """Reject Office zip containers with unsafe expansion characteristics."""
    with zipfile.ZipFile(path) as archive:
        entries = archive.infolist()
        if len(entries) > _MAX_OFFICE_ARCHIVE_ENTRIES:
            raise ValueError("Office document contains too many archive entries")
        total_compressed = sum(max(0, entry.compress_size) for entry in entries)
        total_uncompressed = sum(max(0, entry.file_size) for entry in entries)
        if total_uncompressed > _MAX_OFFICE_UNCOMPRESSED_BYTES:
            raise ValueError("Office document expands beyond the safe limit")
        if (
            total_compressed > 0
            and total_uncompressed / total_compressed
            > _MAX_OFFICE_COMPRESSION_RATIO
        ):
            raise ValueError("Office document compression ratio is unsafe")


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(path)
        if len(reader.pages) > _MAX_PDF_PAGES:
            raise ValueError(f"PDF exceeds the {_MAX_PDF_PAGES}-page extraction limit")
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)
    if suffix == ".docx":
        from docx import Document

        _validate_office_archive(path)
        document = Document(path)
        if len(document.paragraphs) > _MAX_DOCUMENT_PARAGRAPHS:
            raise ValueError("Word document contains too many paragraphs")
        return "\n".join(paragraph.text for paragraph in document.paragraphs)
    if suffix == ".pptx":
        from pptx import Presentation

        _validate_office_archive(path)
        presentation = Presentation(path)
        if len(presentation.slides) > _MAX_PRESENTATION_SLIDES:
            raise ValueError("Presentation contains too many slides")
        values: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            values.append(f"Slide {index}")
            values.extend(
                shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text
            )
        return "\n".join(values)
    if suffix == ".xlsx":
        from openpyxl import load_workbook

        _validate_office_archive(path)
        workbook = load_workbook(path, read_only=True, data_only=True)
        output = io.StringIO()
        writer = csv.writer(output)
        cell_count = 0
        for sheet in workbook.worksheets:
            writer.writerow([f"Sheet: {sheet.title}"])
            for row in sheet.iter_rows(values_only=True):
                cell_count += len(row)
                if cell_count > _MAX_SPREADSHEET_CELLS:
                    raise ValueError("Spreadsheet exceeds the safe cell limit")
                writer.writerow(["" if value is None else value for value in row])
        return output.getvalue()
    return ""


def extract_source_text_from_bytes(name: str, content: bytes) -> str:
    """Extract readable text from in-memory source bytes.

    Unlike ``_extract_text_from_bytes``, extraction failures propagate so course
    sync can record the per-file error instead of silently storing empty text.
    """
    suffix = Path(name).suffix.lower() or ".bin"
    if suffix in TEXT_SUFFIXES:
        return content.decode("utf-8", errors="replace")
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as handle:
        handle.write(content)
        temp_path = Path(handle.name)
    try:
        return extract_text(temp_path)
    finally:
        temp_path.unlink(missing_ok=True)


def _extract_text_from_bytes(name: str, content: bytes) -> str:
    """Extract readable text from in-memory upload bytes for compression checks."""
    suffix = Path(name).suffix.lower() or ".bin"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=True) as handle:
        handle.write(content)
        handle.flush()
        try:
            return extract_text(Path(handle.name))
        except Exception:
            return ""


def _compress_image_bytes(name: str, content: bytes) -> bytes:
    """Downscale and recompress raster images when that shrinks the payload."""
    try:
        from PIL import Image
    except ImportError:
        return content
    try:
        image = Image.open(io.BytesIO(content))
    except Exception:
        return content
    has_alpha = "A" in image.mode or (
        image.mode == "P" and "transparency" in image.info
    )
    width, height = image.size
    scale = min(1.0, _IMAGE_MAX_EDGE / max(width, height, 1))
    if scale < 1.0:
        image = image.resize(
            (max(1, int(width * scale)), max(1, int(height * scale))),
            Image.Resampling.LANCZOS,
        )
    suffix = Path(name).suffix.lower()
    try:
        buffer = io.BytesIO()
        if suffix in {".jpg", ".jpeg"}:
            if image.mode not in {"RGB", "L"}:
                image = image.convert("RGB")
            image.save(
                buffer,
                format="JPEG",
                quality=_IMAGE_JPEG_QUALITY,
                optimize=True,
            )
            candidate = buffer.getvalue()
            return candidate if len(candidate) < len(content) else content
        if suffix == ".webp":
            save_image = image.convert("RGBA") if has_alpha else image.convert("RGB")
            save_image.save(
                buffer,
                format="WEBP",
                quality=_IMAGE_JPEG_QUALITY,
                method=6,
            )
            candidate = buffer.getvalue()
            return candidate if len(candidate) < len(content) else content
        if suffix == ".png":
            save_image = image.convert("RGBA") if has_alpha else image.convert("RGB")
            save_image.save(buffer, format="PNG", optimize=True)
            candidate = buffer.getvalue()
            return candidate if len(candidate) < len(content) else content
        return content
    except Exception:
        return content


def _compress_pdf_bytes(content: bytes) -> bytes:
    """Recompress PDF images when MuPDF is available and text extraction holds."""
    if len(content) < _PDF_COMPRESS_MIN_BYTES:
        return content
    try:
        import fitz
        from PIL import Image
    except ImportError:
        return content

    original_text = _extract_text_from_bytes("source.pdf", content)
    original_chars = len(original_text.strip())
    try:
        document = fitz.open(stream=content, filetype="pdf")
    except Exception:
        return content

    seen: set[int] = set()
    try:
        for page in document:
            for image in page.get_images(full=True):
                xref = int(image[0])
                if xref in seen:
                    continue
                seen.add(xref)
                try:
                    pixmap = fitz.Pixmap(document, xref)
                except Exception:
                    continue
                if pixmap.n - pixmap.alpha >= 4:
                    pixmap = fitz.Pixmap(fitz.csRGB, pixmap)
                if pixmap.alpha:
                    pixmap = fitz.Pixmap(pixmap, 0)
                try:
                    pil_image = Image.frombytes(
                        "RGB",
                        (pixmap.width, pixmap.height),
                        pixmap.samples,
                    )
                except Exception:
                    continue
                width, height = pil_image.size
                scale = min(1.0, _PDF_IMAGE_MAX_EDGE / max(width, height, 1))
                if scale >= 1.0 and width * height < 900_000:
                    continue
                if scale < 1.0:
                    pil_image = pil_image.resize(
                        (max(1, int(width * scale)), max(1, int(height * scale))),
                        Image.Resampling.LANCZOS,
                    )
                buffer = io.BytesIO()
                pil_image.save(
                    buffer,
                    format="JPEG",
                    quality=_PDF_JPEG_QUALITY,
                    optimize=True,
                )
                try:
                    page.replace_image(xref, stream=buffer.getvalue())
                except Exception:
                    continue
        compressed = document.tobytes(
            garbage=4,
            deflate=True,
            clean=True,
            use_objstms=1,
        )
    except Exception:
        logger.debug("PDF compression skipped", exc_info=True)
        return content
    finally:
        try:
            document.close()
        except Exception:
            pass

    if len(compressed) >= len(content):
        return content
    compressed_chars = len(_extract_text_from_bytes("source.pdf", compressed).strip())
    if original_chars and compressed_chars < max(40, int(original_chars * 0.9)):
        return content
    return compressed


def compress_upload_bytes(name: str, content: bytes) -> bytes:
    """Return a smaller payload when safe; otherwise return ``content`` unchanged.

    Used by student-facing upload paths that opt into compression via
    ``save_uploads(..., compress=True)``. Text-bearing PDFs are only replaced
    when extraction quality remains within tolerance. Unsupported or
    already-small files pass through.
    """
    if not content:
        return content
    suffix = Path(name).suffix.lower()
    try:
        if suffix == ".pdf":
            return _compress_pdf_bytes(content)
        if suffix in IMAGE_SUFFIXES and suffix != ".gif":
            return _compress_image_bytes(name, content)
    except Exception:
        logger.debug("Upload compression failed for %s", name, exc_info=True)
    return content


def save_uploads(
    thread_id: str,
    uploads: Iterable[tuple[str, bytes, str | None]],
    *,
    max_file_size_mb: int | None = None,
    compress: bool = True,
    owner_id: str = "local-student",
    source_ids: Sequence[str] | None = None,
) -> list[StoredUpload]:
    """Validate, optionally compress, and store uploads.

    Uses the student-upload size limit by default. When ``compress`` is True,
    compression runs after the size check against the original bytes so
    oversized files are still rejected before any expensive rewrite work.

    Pass ``compress=False`` for trusted copies (for example already-prepared
    lecture-note sync) so notebook create does not re-encode large PDFs.

    Local development keeps the existing ``files_dir/threads/.../uploads``
    layout. Object storage stores bytes under generated keys that include a
    server-generated ``source_id``:

    ``users/<user-id>/notebooks/<notebook-id>/sources/<source-id>/raw/<safe-filename>``
    """
    items = list(uploads)
    if len(items) > settings.max_files:
        raise ValueError(f"Upload at most {settings.max_files} files per message.")
    if source_ids is not None and len(source_ids) != len(items):
        raise ValueError("source_ids length must match uploads")
    size_limit_mb = max_file_size_mb or settings.max_file_size_mb
    for name, content, _ in items:
        if len(content) > size_limit_mb * 1024 * 1024:
            raise ValueError(f"{name} exceeds the {size_limit_mb} MB limit.")
    use_object_storage = settings.file_storage_provider != "local"
    root = None if use_object_storage else _upload_root(thread_id)
    stored: list[StoredUpload] = []
    pending_object_keys: list[str] = []
    try:
        for index, (name, content, supplied_mime) in enumerate(items):
            payload = compress_upload_bytes(name, content) if compress else content
            safe_name = _safe_name(name)
            mime = (
                supplied_mime
                or mimetypes.guess_type(safe_name)[0]
                or "application/octet-stream"
            )
            suffix = Path(safe_name).suffix.lower()
            supported = suffix in SUPPORTED_SUFFIXES
            extracted = ""
            if supported and suffix not in IMAGE_SUFFIXES:
                try:
                    extracted = _extract_text_from_bytes(safe_name, payload)
                except Exception as exc:
                    extracted = f"[Could not extract {safe_name}: {type(exc).__name__}]"

            if use_object_storage:
                from backend.persistence.factory import get_file_storage
                from backend.persistence.object_keys import build_upload_object_key

                source_id = str(
                    (source_ids[index] if source_ids is not None else None)
                    or uuid.uuid4()
                )
                key = build_upload_object_key(
                    user_id=owner_id,
                    notebook_id=thread_id,
                    source_id=source_id,
                    filename=safe_name,
                )
                # Track before PUT: some clients can create the object and then
                # raise while decoding a response. Deleting a missing key is safe.
                pending_object_keys.append(key)
                get_file_storage().put_bytes(key=key, data=payload, content_type=mime)
                temp_path = Path(tempfile.gettempdir()) / safe_name
                stored.append(
                    StoredUpload(
                        name=safe_name,
                        path=temp_path,
                        mime=mime,
                        size=len(payload),
                        supported=supported,
                        extracted_text=extracted[:120_000],
                        storage_key=key,
                        storage_provider=settings.file_storage_provider,
                        source_id=source_id,
                    )
                )
                continue

            assert root is not None
            candidate = root / safe_name
            counter = 2
            while candidate.exists():
                candidate = root / f"{Path(safe_name).stem}-{counter}{Path(safe_name).suffix}"
                counter += 1
            candidate.write_bytes(payload)
            if supported and suffix not in IMAGE_SUFFIXES and not extracted:
                try:
                    extracted = extract_text(candidate)
                except Exception as exc:
                    extracted = (
                        f"[Could not extract {candidate.name}: {type(exc).__name__}]"
                    )
            local_source_id = (
                str(source_ids[index]) if source_ids is not None else None
            )
            stored.append(
                StoredUpload(
                    name=candidate.name,
                    path=candidate,
                    mime=mime,
                    size=len(payload),
                    supported=supported,
                    extracted_text=extracted[:120_000],
                    storage_provider="local",
                    source_id=local_source_id,
                )
            )
    except Exception:
        if pending_object_keys:
            from backend.persistence.factory import get_file_storage

            storage = get_file_storage()
            for key in pending_object_keys:
                try:
                    storage.delete(key)
                except Exception:
                    logger.warning(
                        "Failed to clean object after upload batch failure",
                        exc_info=True,
                    )
        raise
    return stored


def image_input(upload: StoredUpload) -> dict[str, str]:
    encoded = base64.b64encode(upload.read_bytes()).decode("ascii")
    return {
        "type": "input_image",
        "image_url": f"data:{upload.mime};base64,{encoded}",
        "detail": "auto",
    }


def document_context(uploads: Iterable[StoredUpload], limit: int = 160_000) -> str:
    sections: list[str] = []
    remaining = limit
    for upload in uploads:
        if not upload.extracted_text:
            continue
        text = upload.extracted_text[:remaining]
        sections.append(f"--- {upload.name} ---\n{text}")
        remaining -= len(text)
        if remaining <= 0:
            break
    return "\n\n".join(sections)
